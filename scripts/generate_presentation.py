import sys
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_presentation(output_path="ZamRoam_Tourism_Platform_Presentation.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Brand Colors
    C_DARK_BG = RGBColor(13, 43, 39)       # #0D2B27 Deep Teal
    C_CARD_BG = RGBColor(18, 56, 51)       # #123833 Card Teal
    C_COPPER = RGBColor(222, 119, 57)      # #DE7739 Copper/Orange
    C_GOLD = RGBColor(245, 158, 11)        # #F59E0B Gold
    C_EMERALD = RGBColor(52, 211, 153)     # #34D399 Bright Emerald
    C_CYAN = RGBColor(56, 189, 248)        # #38BDF8 Sky Cyan
    C_WHITE = RGBColor(255, 255, 255)
    C_MUTED = RGBColor(163, 207, 201)      # #A3CFC9 Light Muted Mint

    def create_slide_background(slide, color=C_DARK_BG):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_header(slide, category, title, subtitle=None):
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p0 = tf.paragraphs[0]
        p0.text = category.upper()
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.color.rgb = C_COPPER
        
        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.size = Pt(22)
        p1.font.bold = True
        p1.font.color.rgb = C_WHITE
        p1.space_before = Pt(3)

        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(11)
            p2.font.color.rgb = C_MUTED
            p2.space_before = Pt(2)

    def add_card(slide, left, top, width, height, title, items, badge=None, border_color=C_COPPER, bg_color=C_CARD_BG):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        
        tb = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.2), Inches(width - 0.5), Inches(height - 0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        if badge:
            p_badge = tf.paragraphs[0]
            p_badge.text = badge.upper()
            p_badge.font.size = Pt(9)
            p_badge.font.bold = True
            p_badge.font.color.rgb = C_COPPER
            p_title = tf.add_paragraph()
        else:
            p_title = tf.paragraphs[0]
            
        p_title.text = title
        p_title.font.size = Pt(14)
        p_title.font.bold = True
        p_title.font.color.rgb = C_WHITE
        p_title.space_after = Pt(8)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(10.5)
            p.font.color.rgb = C_MUTED
            p.space_after = Pt(4)

    # 1. Title Slide
    s1 = prs.slides.add_slide(blank_slide_layout)
    create_slide_background(s1, C_DARK_BG)
    accent = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.12), Inches(4.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = C_COPPER
    accent.line.fill.background()

    tb = s1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(11.0), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ZAMROAM"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = C_WHITE

    p2 = tf.add_paragraph()
    p2.text = "Roam Zambia. Experience More."
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = C_COPPER
    p2.space_before = Pt(4)

    p3 = tf.add_paragraph()
    p3.text = "The National Tourism Technology & Digital Travel Ecosystem Connecting 10 Provinces, 116 Districts, National Parks, Traditional Ceremonies, and Verified Local Providers."
    p3.font.size = Pt(13)
    p3.font.color.rgb = C_MUTED
    p3.space_before = Pt(18)

    p4 = tf.add_paragraph()
    p4.text = "🇿🇲 Commercial Tourism Platform Owned & Operated by Lamton Investments Ltd\n📍 Lusaka • Livingstone • Ndola • Mfuwe  |  🌐 zamroam.com  |  💬 WhatsApp: +260 573 506 598"
    p4.font.size = Pt(11)
    p4.font.color.rgb = C_EMERALD
    p4.space_before = Pt(24)

    # 2. Executive Overview
    s2 = prs.slides.add_slide(blank_slide_layout)
    create_slide_background(s2)
    add_header(s2, "Executive Overview", "Transforming Zambia Tourism with Modern Digital Infrastructure", "Connecting international travellers and domestic explorers directly with verified local experiences.")
    add_card(s2, 0.8, 1.8, 3.6, 5.0, "The Market Challenge", [
        "Fragmented Tourism Discovery: Safari lodges, bush camps, and local guides lacked unified digital presence.",
        "Offline & Remote Barriers: Spotty cellular connectivity in national parks hindered digital bookings & navigation.",
        "Middleman Commissions: High international aggregator fees (20-30%) reduced local operator revenues.",
        "Logistical Complexity: Difficult travel time estimates and lack of en-route amenity data across 10 provinces."
    ], badge="Problem Statement", border_color=RGBColor(239, 68, 68))
    add_card(s2, 4.8, 1.8, 3.6, 5.0, "The ZamRoam Solution", [
        "Unified 10-Province Hub: Complete digital coverage of 116 districts, national parks, and cultural heritage sites.",
        "Offline-First Native Architecture: Topographic trail maps, GPX waypoints, and passes accessible without cellular signal.",
        "100 Founding Partners Program: Zero-commission model empowering Zambian lodges, operators, and guides.",
        "Smart GIS & Slippy Tiles: Interactive OpenStreetMap with live highway corridor travel time & distance calculations."
    ], badge="Core Value Proposition", border_color=C_EMERALD)
    add_card(s2, 8.8, 1.8, 3.7, 5.0, "Platform Key Metrics", [
        "10 Provinces & 116 Districts with precise administrative & provincial capital headquarters mapping.",
        "11 National Traditional Ceremonies with authentic cultural guides and 7-language phrasebooks.",
        "8 Major National Highway Corridors (T1, T2, T3, T4, T5, M3, M9, M10) with realistic road waypoints.",
        "24/7 Direct Emergency Dispatch (Zambia Police, National Ambulance, Medevac, ZTA Tourism Police)."
    ], badge="Key Capabilities", border_color=C_GOLD)

    # 3. GIS & Map
    s3 = prs.slides.add_slide(blank_slide_layout)
    create_slide_background(s3)
    add_header(s3, "GIS & Navigation Engine", "Enterprise OpenStreetMap Slippy Tiles & Road Corridor Routing", "Dynamic travel calculations with realistic highway geometry instead of straight-line vectors.")
    add_card(s3, 0.8, 1.8, 3.6, 5.0, "Dynamic Distance & TTT", [
        "Lusaka National HQ Hub: Instant road distance & travel time estimation from the capital city.",
        "Provincial Capital Hubs: Strict mapping for all 10 provincial HQs (e.g. Solwezi for North-Western, Mongu for Western, Kabwe for Central).",
        "District Headquarters: Live multi-point distance breakdown for the host district.",
        "4 Travel Modes: Sedan (85 km/h), Overland 4x4 (65 km/h), Long-Distance Bus (60 km/h), and Air Charter (280 km/h)."
    ], badge="Logistics Calculator", border_color=C_CYAN)
    add_card(s3, 4.8, 1.8, 3.6, 5.0, "Real Highway Corridors", [
        "T1 Tourism Highway: Livingstone / Victoria Falls – Choma – Mazabuka – Kafue – Lusaka.",
        "T4 Great East Corridor: Lusaka – Chongwe – Luangwa Bridge – Petauke – Chipata – Mfuwe.",
        "T5 North-Western Highway: Chingola – Solwezi – Mutanda – Mwinilunga – Ikelenge (Zambezi Source).",
        "T2 Great North Road & M9/M10 Western Barotseland Highways with full waypoint curves."
    ], badge="Highway Network", border_color=C_COPPER)
    add_card(s3, 8.8, 1.8, 3.7, 5.0, "Interactive Points of Interest", [
        "Interactive Destination Pins: Highlighting UNESCO sites, national parks, waterfalls, and cultural hubs.",
        "National Airports & Bush Airstrips: Kenneth Kaunda, Harry Mwaanga Nkumbula, Mfuwe, Zengamina, Ndole Bay.",
        "24/7 En-Route Amenities: Fuel stations, Level 1 trauma emergency hospitals, and cellular coverage bands.",
        "Smooth Mouse Wheel & Double-Click Zoom with instant re-projection."
    ], badge="POI Layering", border_color=C_EMERALD)

    # 4. Catalogue & Culture
    s4 = prs.slides.add_slide(blank_slide_layout)
    create_slide_background(s4)
    add_header(s4, "Tourism & Cultural Assets", "Curated Zambian Experiences & Authentic Living Heritage", "Empowering international tourists and domestic travellers to explore Zambia beyond conventional routes.")
    add_card(s4, 0.8, 1.8, 3.6, 5.0, "7 Core Listing Categories", [
        "Safari Lodges & Stays: Luxury riverfront chalets, eco-camps, and boutique city hotels.",
        "Safaris & Guided Tours: Birthplace of walking safaris in South Luangwa, canoe trails in Lower Zambezi.",
        "Adventure & Extreme: Devil's Pool, Victoria Falls bungee, gorge swings, microlight flights.",
        "Dining & Nightlife: Traditional Zambian boma feasts, fresh Kariba bream, and craft breweries.",
        "Transport & Overland 4x4 Car Hire with GPS kits."
    ], badge="Catalogue Breadth", border_color=C_GOLD)
    add_card(s4, 4.8, 1.8, 3.6, 5.0, "11 Traditional Ceremonies", [
        "Kuomboka (Western / Barotseland): Royal barge Nalikwanda voyage from Lealui to Limulunga.",
        "Nc'wala (Eastern / Chipata): Ngoni first-fruits harvest celebration and warrior dances.",
        "Likumbi Lya Mize (North-Western / Zambezi): UNESCO-inscribed Makishi masquerade.",
        "Ukusefya Pa Ng'wena, Mutomboko, Shimunenga, Kulamba, Lwiindi Gonde, and more."
    ], badge="Living Culture", border_color=C_COPPER)
    add_card(s4, 8.8, 1.8, 3.7, 5.0, "7-Language Phrasebook", [
        "Comprehensive phrasebook engine with audio-ready greetings, transport, bargaining, and emergency phrases.",
        "Languages Supported: Bemba, Nyanja (Chewa), Tonga, Lozi, Lunda, Luvale, and Kaonde.",
        "Localized Cultural Etiquette: Royal court greetings, chieftaincy respect, and regional customs.",
        "100% accessible offline without data roaming."
    ], badge="Language Engine", border_color=C_EMERALD)

    # 5. Pass & Membership
    s5 = prs.slides.add_slide(blank_slide_layout)
    create_slide_background(s5)
    add_header(s5, "Commercial Engine", "The ZamRoam Pass & Digital Membership Loyalty Ecosystem", "Tiered digital memberships unlocking verified provider discounts, VIP access, and anti-fraud QR passes.")
    add_card(s5, 0.8, 1.8, 3.6, 5.0, "Tiered Membership Passes", [
        "Tourist Explorer Pass (Complimentary): Basic traveler profile, wishlist saving, review moderation.",
        "Safari Adventurer Pass (K499/year or K49/mo): 10-15% discounts at partner lodges, walking safaris, priority booking.",
        "VIP Diplomat Pass (K1,699/year): Exclusive room upgrades, 4 family dependants, private VIP airport transfers.",
        "7-Day Event & Conference Pass (K149): Short-term pass for delegates and holidaymakers."
    ], badge="Membership Tiers", border_color=C_EMERALD)
    add_card(s5, 4.8, 1.8, 3.6, 5.0, "Dynamic QR Tokens", [
        "Time-Expiring Security Salt: Anti-fraud dynamic QR codes refreshing every 60 seconds to prevent duplication.",
        "Offline-Ready Wallet: Passes cached in local storage with cryptographic verification tokens.",
        "Physical NFC/Metal Cards: Premium engraved metal membership cards dispatched to VIP members.",
        "Real-Time Status: Instant validation across mobile, tablet, and desktop devices."
    ], badge="Digital Security", border_color=C_CYAN)
    add_card(s5, 8.8, 1.8, 3.7, 5.0, "Gamified Loyalty Ledger", [
        "Reward Points Engine: Earn points on verified reviews, partner visits, and multi-day bookings.",
        "Redemption Catalog: Convert points to park activity vouchers, lodge discounts, or dining credits.",
        "Provider Redemption Terminal: Web-based scanning portal for partner lodge staff to validate passes in seconds.",
        "Complete audit logging of all redemptions."
    ], badge="Loyalty Engine", border_color=C_GOLD)

    # 6. Founding Partners
    s6 = prs.slides.add_slide(blank_slide_layout)
    create_slide_background(s6)
    add_header(s6, "B2B Strategy & Growth", "100 Founding Partners Program & Sustainable Revenue Model", "Empowering Zambian tourism businesses while establishing diversified recurring revenue streams.")
    add_card(s6, 0.8, 1.8, 3.6, 5.0, "100 Founding Partners", [
        "0% Commission Year 1: Zero listing and transaction fees for the first 100 verified Zambian tourism businesses.",
        "Founding Partner Badge: Permanent platform distinction establishing verified trust and credibility.",
        "Dedicated Discovery Placement: Featured spotlights across search, interactive map, and editorial circuits.",
        "Direct Booking Leads: Customers connect directly with provider WhatsApp, website, and phone."
    ], badge="Partner Incentive", border_color=C_COPPER)
    add_card(s6, 4.8, 1.8, 3.6, 5.0, "Provider Subscriptions", [
        "Listed Partner (Free): Verified directory listing, up to 2 active member offers.",
        "Silver Partner (K899/yr): Enhanced search ranking, 4 simultaneous member offers, analytics dashboard.",
        "Gold Partner (K1,899/yr): Priority discovery, geographic push targeting, 8 active member offers.",
        "Platinum Partner (K4,500/yr): Homepage spotlight campaigns, VIP concierge integration, API readiness."
    ], badge="B2B Revenue", border_color=C_GOLD)
    add_card(s6, 8.8, 1.8, 3.7, 5.0, "Diversified Income Streams", [
        "Annual Consumer Pass Subscriptions (B2C recurring membership revenue).",
        "Corporate & Diplomatic Sponsorships (custom hospitality packages for corporate teams).",
        "Featured Partner Advertising & Regional Circuit Sponsorships.",
        "Transactional Commission Options for automated high-volume bookings."
    ], badge="Business Model", border_color=C_EMERALD)

    # 7. AI & Safety
    s7 = prs.slides.add_slide(blank_slide_layout)
    create_slide_background(s7)
    add_header(s7, "AI & Safety Architecture", "AI Travel Concierge, Security Advisory & Emergency Dispatch", "State-of-the-art traveler assistance and comprehensive multi-agency safety infrastructure.")
    add_card(s7, 0.8, 1.8, 3.6, 5.0, "AI Travel Concierge", [
        "24/7 Smart Assistance: Context-aware AI assistant providing instant answers on visas, health, and transport.",
        "Dynamic Itinerary Builder: Generates custom 3-day, 7-day, or 14-day safari routes tailored to budget and travel style.",
        "Multi-Channel Support: Interactive in-app chat modal with seamless escalation to human WhatsApp concierge.",
        "Offline fallback guidance for remote areas."
    ], badge="AI Technology", border_color=C_CYAN)
    add_card(s7, 4.8, 1.8, 3.6, 5.0, "Security Advisory Engine", [
        "Provincial Safety Indexes: Real-time safety ratings, weather advisories, and road condition updates.",
        "Verified Local Advisories: National park flood warnings, dry-season wildlife movements, and border crossing tips.",
        "Health Guidance: Malaria prophylaxis advice, yellow fever requirements, and accredited clinic directories.",
        "Transparent, non-alarmist traveler guidance."
    ], badge="Safety Index", border_color=C_GOLD)
    add_card(s7, 8.8, 1.8, 3.7, 5.0, "1-Tap Emergency Dispatch", [
        "Zambia Police Service (National Hotline & Local Stations).",
        "National Emergency Medical Ambulance (Level 1 Trauma Centers).",
        "Medevac Aerial Rescue (Flying Mission Zambia & ProCharter Air Evac).",
        "Zambia Tourism Agency (ZTA) Tourist Police Helpline.",
        "Direct click-to-call integration across all devices."
    ], badge="Emergency Ready", border_color=RGBColor(239, 68, 68))

    # 8. Tech Architecture
    s8 = prs.slides.add_slide(blank_slide_layout)
    create_slide_background(s8)
    add_header(s8, "Technical Architecture", "Enterprise Next.js Architecture, Cloudflare D1 & Strict RBAC", "High-performance, secure, and production-hardened web application infrastructure.")
    add_card(s8, 0.8, 1.8, 3.6, 5.0, "Frontend & UX Engineering", [
        "Next.js 16 (App Router) + React + TypeScript.",
        "Zero Heavy CSS Frameworks: Native responsive CSS with tailored design tokens, glassmorphism, and Ubuntu typography.",
        "Vector SVG Canvas Engine: Smooth GPS projections for 116 districts, highway corridors, and custom interactive pins.",
        "Accessible, touch-friendly, and mobile-optimized."
    ], badge="Frontend Stack", border_color=C_EMERALD)
    add_card(s8, 4.8, 1.8, 3.6, 5.0, "Database & Storage Layer", [
        "Cloudflare D1 / SQLite Engine with ACID compliance and parameterized queries.",
        "Strict Non-Destructive Migrations: Safe additive column upgrades without table drops or wipes.",
        "Multi-Entity Relationships: Countries, Provinces, Districts, Categories, Listings, Reviews, Wishlists, Bookings, Memberships.",
        "Complete audit logging of all sensitive actions."
    ], badge="Data Safety", border_color=C_CYAN)
    add_card(s8, 8.8, 1.8, 3.7, 5.0, "Security & Admin Control", [
        "Multi-Tier Role-Based Access Control (RBAC): SuperAdmin, Regional Moderator, Provider Owner, Traveler.",
        "Provider Moderation Portal: Review submitted listings, license credentials, and promotional offers before publication.",
        "Anti-Abuse & Rate Limiting: Secure auth token validation with HTTP-only cookie persistence.",
        "Automated Test Coverage (26/26 test suites)."
    ], badge="Enterprise RBAC", border_color=C_COPPER)

    # 9. Roadmap
    s9 = prs.slides.add_slide(blank_slide_layout)
    create_slide_background(s9)
    add_header(s9, "Vision & Roadmap", "National Scaling, Regional Expansion & Long-Term Impact", "Positioning Zambia at the forefront of digital tourism innovation across Africa.")
    add_card(s9, 0.8, 1.8, 3.6, 5.0, "Phase 1: National Launch", [
        "Complete onboarding of the 100 Founding Zambian Tourism Partners.",
        "Nationwide roll-out of the ZamRoam Pass across Lusaka, Livingstone, Ndola, and South Luangwa.",
        "Partnership alignment with ZTA (Zambia Tourism Agency) and NHCC (National Heritage).",
        "Initial domestic & international tourist adoption."
    ], badge="Q1 - Q2 2026", border_color=C_EMERALD)
    add_card(s9, 4.8, 1.8, 3.6, 5.0, "Phase 2: Ecosystem Growth", [
        "Integration with domestic airline reservation APIs (Zambia Airways & Proflight).",
        "Mobile POS & QR scanner hardware distribution to partner lodges and national park gates.",
        "Offline PWA (Progressive Web App) with background synchronization and trail telemetry.",
        "Expansion of corporate membership pass sales."
    ], badge="Q3 - Q4 2026", border_color=C_GOLD)
    add_card(s9, 8.8, 1.8, 3.7, 5.0, "Phase 3: SADC Expansion", [
        "Cross-border safari circuits connecting Zambia, Botswana (Chobe), Zimbabwe (Vic Falls), and Namibia (Caprivi).",
        "Multi-currency automated settlement (ZMW, USD, EUR, GBP, ZAR).",
        "Advanced AI dynamic pricing & demand forecasting for local operators.",
        "Regional tourism data analytics engine."
    ], badge="2027 & Beyond", border_color=C_COPPER)

    # 10. Conclusion & Contact
    s10 = prs.slides.add_slide(blank_slide_layout)
    create_slide_background(s10, C_DARK_BG)
    accent10 = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(0.12), Inches(5.1))
    accent10.fill.solid()
    accent10.fill.fore_color.rgb = C_COPPER
    accent10.line.fill.background()

    tb10 = s10.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(11.0), Inches(5.1))
    tf10 = tb10.text_frame
    tf10.word_wrap = True

    p_c1 = tf10.paragraphs[0]
    p_c1.text = "Partner With ZamRoam"
    p_c1.font.size = Pt(36)
    p_c1.font.bold = True
    p_c1.font.color.rgb = C_WHITE

    p_c2 = tf10.add_paragraph()
    p_c2.text = "Join Zambia's Digital Tourism Revolution"
    p_c2.font.size = Pt(20)
    p_c2.font.bold = True
    p_c2.font.color.rgb = C_COPPER
    p_c2.space_before = Pt(4)

    p_c3 = tf10.add_paragraph()
    p_c3.text = "ZamRoam is ready to partner with lodge operators, safari guides, transport providers, cultural institutions, and government stakeholders to showcase the best of Zambia to the world."
    p_c3.font.size = Pt(13)
    p_c3.font.color.rgb = C_MUTED
    p_c3.space_before = Pt(16)

    p_c4 = tf10.add_paragraph()
    p_c4.text = "CORPORATE & PARTNERSHIP INQUIRIES:"
    p_c4.font.size = Pt(11)
    p_c4.font.bold = True
    p_c4.font.color.rgb = C_GOLD
    p_c4.space_before = Pt(22)

    p_c5 = tf10.add_paragraph()
    p_c5.text = "🏢 Legal Entity: Lamton Investments Ltd\n📍 Registered Office: Plot 10444, Great East Road, Rhodes Park, Lusaka, Zambia\n🌐 Official Website: https://zamroam.com\n✉️ General Inquiries: info@zamroam.com  |  info@lamtoninvestments.com\n💬 WhatsApp Business Hotline: +260 573 506 598"
    p_c5.font.size = Pt(12)
    p_c5.font.color.rgb = C_WHITE
    p_c5.space_before = Pt(6)

    prs.save(output_path)
    print(f"Successfully generated PowerPoint presentation at {output_path}")

if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "ZamRoam_Tourism_Platform_Presentation.pptx"
    build_presentation(out)
